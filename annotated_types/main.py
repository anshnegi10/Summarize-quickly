from dotenv import load_dotenv
load_dotenv()
import os
from datetime import datetime, timedelta, timezone
import json
import pandas as pd
import requests
import botocore
import boto3

# ------------------ AWS CONFIG ------------------- #
REGION = os.getenv('AWS_REGION', 'us-east-1')
S3_BUCKET = os.getenv('S3_BUCKET', 'smart-notes-uploads')
DYNAMO_TABLE = os.getenv('DYNAMO_TABLE', '')  # Set a table name to activate DB logging
TOGETHER_API_KEY = os.getenv('TOGETHER_API_KEY', '')
# ------------------------------------------------- #

# AWS Clients
cloudwatch = boto3.client('cloudwatch', region_name=REGION)
ec2 = boto3.client('ec2', region_name=REGION)
pricing = boto3.client('pricing', region_name=REGION)
s3 = boto3.client('s3', region_name=REGION)
dynamodb = boto3.resource('dynamodb', region_name=REGION) if DYNAMO_TABLE else None

# --- EC2 COST OPTIMIZATION ----

def get_instance_metrics(instance_id, days=14):
    end_time = datetime.now(timezone.utc)
    start_time = end_time - timedelta(days=days)
    response = cloudwatch.get_metric_statistics(
        Namespace='AWS/EC2',
        MetricName='CPUUtilization',
        Dimensions=[{'Name': 'InstanceId', 'Value': instance_id}],
        StartTime=start_time,
        EndTime=end_time,
        Period=86400,
        Statistics=['Average', 'Maximum'],
        Unit='Percent'
    )
    datapoints = response.get('Datapoints', [])
    if not datapoints:
        return None, None
    avg_util = sum(dp['Average'] for dp in datapoints) / len(datapoints)
    max_util = max(dp['Maximum'] for dp in datapoints)
    return avg_util, max_util

def get_all_instances():
    response = ec2.describe_instances(Filters=[{'Name': 'instance-state-name', 'Values': ['running']}])
    instances = []
    for reservation in response['Reservations']:
        for instance in reservation['Instances']:
            instances.append({
                'InstanceId': instance['InstanceId'],
                'InstanceType': instance['InstanceType'],
                'LaunchTime': instance['LaunchTime'],
                'Tags': {tag['Key']: tag['Value'] for tag in instance.get('Tags', [])}
            })
    return instances

def get_instance_type_info(instance_type):
    try:
        response = pricing.get_products(
            ServiceCode='AmazonEC2',
            Filters=[
                {'Type': 'TERM_MATCH', 'Field': 'instanceType', 'Value': instance_type},
                {'Type': 'TERM_MATCH', 'Field': 'location', 'Value': 'US East (N. Virginia)'},
                {'Type': 'TERM_MATCH', 'Field': 'operatingSystem', 'Value': 'Linux'},
                {'Type': 'TERM_MATCH', 'Field': 'preInstalledSw', 'Value': 'NA'},
                {'Type': 'TERM_MATCH', 'Field': 'tenancy', 'Value': 'Shared'},
                {'Type': 'TERM_MATCH', 'Field': 'capacitystatus', 'Value': 'Used'}
            ],
            MaxResults=1
        )
        if not response['PriceList']:
            return None
        price_item = json.loads(response['PriceList'][0])
        terms = price_item['terms']['OnDemand']
        term_keys = list(terms.keys())
        price_dimensions = terms[term_keys[0]]['priceDimensions']
        dimension_keys = list(price_dimensions.keys())
        usd_price = price_dimensions[dimension_keys[0]]['pricePerUnit'].get('USD', None)
        return float(usd_price) if usd_price else None
    except Exception as e:
        print(f"Error getting pricing for {instance_type}: {e}")
        return None

def get_recommendations(instance, avg_util, max_util):
    recommendations = []
    current_type = instance['InstanceType']
    if max_util < 40:
        recommendations.append(f"Underutilized (Max CPU: {max_util:.1f}%)")
        current_price = get_instance_type_info(current_type)
        family = current_type.split('.')[0]
        potential_types = [f"{family}.large", f"{family}.medium", f"{family}.small"]
        for new_type in potential_types:
            if new_type == current_type:
                continue
            new_price = get_instance_type_info(new_type)
            if new_price and current_price and new_price < current_price:
                savings = current_price - new_price
                savings_percent = (savings / current_price) * 100
                recommendations.append(
                    f"⬇️ Suggest: {new_type} → Save ${savings:.2f}/hr ({savings_percent:.1f}%)"
                )
    if max_util < 10:
        recommendations.append("Consider stopping: appears unused")
    return recommendations

def upload_to_s3(file_name, bucket_name, object_name=None):
    if object_name is None:
        object_name = os.path.basename(file_name)
    try:
        s3.upload_file(file_name, bucket_name, object_name)
        print(f"✅ Uploaded to s3://{bucket_name}/{object_name}")
        return True
    except botocore.exceptions.ClientError as e:
        print(f"❌ Failed to upload to S3: {e}")
        return False

def write_to_dynamodb(table_name, record):
    if not dynamodb:
        print("⚠️ DynamoDB not configured.")
        return
    table = dynamodb.Table(table_name)
    table.put_item(Item=record)
    print(f"✅ Record logged to DynamoDB table: {table_name}")

def analyze_instances(do_shutdown=None, do_start=None, dynamo_logging=False):
    print("\n AWS EC2 Instance Cost Optimization Tool")
    print("=" * 60)
    instances = get_all_instances()
    print(f"Found {len(instances)} running EC2 instance(s)\n")
    results = []
    for instance in instances:
        instance_id = instance['InstanceId']
        print(f" Analyzing {instance_id}...")
        try:
            avg_util, max_util = get_instance_metrics(instance_id)
            if avg_util is None:
                print(" No metrics found.")
                continue
            recommendations = get_recommendations(instance, avg_util, max_util)
            results.append({
                'InstanceId': instance_id,
                'Name': instance['Tags'].get('Name', ''),
                'InstanceType': instance['InstanceType'],
                'AvgCPU': round(avg_util, 1),
                'MaxCPU': round(max_util, 1),
                'Recommendations': "\n".join(recommendations) if recommendations else "✅ No recommendations"
            })
        except Exception as e:
            print(f" Error with {instance_id}: {e}")
    df = pd.DataFrame(results)
    if df.empty:
        print("\n No data to export..............")
        return None
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
    file_name = f"aws_cost_optimization_report_{timestamp}.csv"
    full_path = os.path.abspath(file_name)
    df.to_csv(file_name, index=False)
    print("\n📄 Final Report:")
    print(df.to_string(index=False))
    print(f"\n Report saved at: {full_path}")
    uploaded = upload_to_s3(file_name, S3_BUCKET)
    if dynamo_logging and DYNAMO_TABLE:
        for record in results:
            record['created_at'] = datetime.now().isoformat()
            write_to_dynamodb(DYNAMO_TABLE, record)
    # AI Summarize the report!
    ai_summary_to_file(file_name)
    return {"file": file_name, "s3": uploaded, "items": results}

def summarize_with_together(text):
    if not TOGETHER_API_KEY:
        print("No Together API key set; skipping summarization.")
        return ""
    response = requests.post(
        "https://api.together.xyz/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {TOGETHER_API_KEY}",
            "Content-Type": "application/json"
        },
        json={
            "model": "meta-llama/Llama-3-8b-chat-hf",
            "messages": [{"role": "user", "content": "Summarize this document:\n" + text[:3000]}],
            "max_tokens": 350
        }
    )
    output = response.json()
    return output["choices"][0]["message"]["content"] if "choices" in output else str(output)

def ai_summary_to_file(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()
    summary = summarize_with_together(content)
    if summary.strip():
        summary_path = file_path.replace('.csv', '_summary.txt').replace('.pdf', '_summary.txt').replace('.pptx', '_summary.txt')
        with open(summary_path, "w", encoding="utf-8") as f:
            f.write(summary)
        print("\n====== AI SUMMARY OF FILE ======")
        print(summary)
        print(f"\nSummary also saved at: {summary_path}\n")
        upload_to_s3(summary_path, S3_BUCKET)
        # Log summary to DynamoDB as well (if active)
        if DYNAMO_TABLE:
            write_to_dynamodb(DYNAMO_TABLE, {
                'filename': os.path.basename(summary_path),
                'created_at': datetime.now().isoformat(),
                'summary_preview': summary[:200]
            })
    else:
        print("\n[!] Summary not generated.")

# ----------- PPTX AND PDF SUMMARIZATION ------
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

def extract_text_from_pptx(filepath):
    if not HAS_PPTX:
        raise RuntimeError("python-pptx not installed.")
    prs = Presentation(filepath)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)

def pptx_summarize_handler(file_path: str, output_txt: str = None):
    if not HAS_PPTX:
        print("python-pptx not available.")
        return
    ppt_text = extract_text_from_pptx(file_path)
    summary = summarize_with_together(ppt_text)
    if not output_txt:
        output_txt = file_path.replace('.pptx', '_summary.txt')
    with open(output_txt, "w", encoding="utf-8") as f:
        f.write(summary)
    print(f"PPTX summary written to: {output_txt}")
    upload_to_s3(output_txt, S3_BUCKET)
    if DYNAMO_TABLE:
        write_to_dynamodb(DYNAMO_TABLE, {
            'filename': os.path.basename(file_path),
            'summary_s3_key': os.path.basename(output_txt),
            'created_at': datetime.now().isoformat(),
            'summary_preview': summary[:200]
        })

def pdf_summarize_handler(file_path: str, output_txt: str = None):
    import pdfplumber
    with pdfplumber.open(file_path) as pdf:
        text = '\n'.join(page.extract_text() or '' for page in pdf.pages)
    summary = summarize_with_together(text)
    if not output_txt:
        output_txt = file_path.replace('.pdf', '_summary.txt')
    with open(output_txt, "w", encoding="utf-8") as f:
        f.write(summary)
    print(f"PDF summary written to: {output_txt}")
    upload_to_s3(output_txt, S3_BUCKET)
    if DYNAMO_TABLE:
        write_to_dynamodb(DYNAMO_TABLE, {
            'filename': os.path.basename(file_path),
            'summary_s3_key': os.path.basename(output_txt),
            'created_at': datetime.now().isoformat(),
            'summary_preview': summary[:200]
        })

# -------- MAIN MENU / CLI ENTRY ----------
if __name__ == "__main__":
    analyze_instances()
    # Summarize local PPTX or PDF
    doc_path = input("\nEnter full path to a PPTX or PDF to summarize (or press Enter to skip): ").strip()
    if doc_path:
        if doc_path.lower().endswith('.pptx'):
            pptx_summarize_handler(doc_path)
        elif doc_path.lower().endswith('.pdf'):
            pdf_summarize_handler(doc_path)
        else:
            print("File type not supported for summarization.")
