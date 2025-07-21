import os
import boto3
import requests

def extract_pptx_text(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return '\n'.join(text)

def extract_pdf_text(file_path):
    import pdfplumber
    with pdfplumber.open(file_path) as pdf:
        return ' '.join(page.extract_text() or '' for page in pdf.pages)

def summarize_with_together(text):
    TOGETHER_API_KEY = os.getenv('TOGETHER_API_KEY')
    response = requests.post(
        "https://api.together.xyz/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {TOGETHER_API_KEY}",
            "Content-Type": "application/json"
        },
        json={
            "model": "meta-llama/Llama-3-8b-chat-hf",
            "messages": [{"role": "user", "content": f"Summarize these notes:\n\n{text[:3000]}"}],
            "max_tokens": 350
        }
    )
    output = response.json()
    try:
        return output["choices"][0]["message"]["content"]
    except Exception:
        return str(output)

def lambda_handler(event, context):
    s3 = boto3.client('s3')
    bucket = event['Records'][0]['s3']['bucket']['name']
    key = event['Records'][0]['s3']['object']['key']
    filename = '/tmp/input_file'
    s3.download_file(bucket, key, filename)

    if key.lower().endswith('.pptx'):
        text = extract_pptx_text(filename)
        summary_filename = key.replace('.pptx', '_summary.txt')
    elif key.lower().endswith('.pdf'):
        text = extract_pdf_text(filename)
        summary_filename = key.replace('.pdf', '_summary.txt')
    else:
        return {'statusCode': 400, 'body': 'Unsupported file type'}

    summary = summarize_with_together(text)
    s3.put_object(Bucket=bucket, Key=summary_filename, Body=summary.encode('utf-8'))
    return {'statusCode': 200, 'body': f'Summary saved as {summary_filename}'}
